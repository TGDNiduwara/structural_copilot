"""
tools/pptx_tool.py
==================
Engineering presentation generator using python-pptx.

Produces a presentation deck (.pptx) with a title slide, project
assumptions, design standards, executive summary, governing member-forces
table, support reactions table, and embedded SFD/BMD diagram images —
visually consistent with the Excel and Word tools (navy / steel-blue
engineering theme).

Requires: python-pptx (`pip install python-pptx`)
"""

from __future__ import annotations

import logging
import os
from datetime import datetime
from typing import List, Optional

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

logger = logging.getLogger("structural_copilot.pptx_tool")
logger.setLevel(logging.INFO)

NAVY = RGBColor(0x1F, 0x38, 0x64)
STEEL_BLUE = RGBColor(0x2E, 0x75, 0xB6)
GREY = RGBColor(0x59, 0x59, 0x59)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x33, 0x33, 0x33)
LIGHT_BLUE = RGBColor(0xD9, 0xE1, 0xF2)

# 16:9 slide geometry (inches)
SLIDE_W = 13.333
SLIDE_H = 7.5
MARGIN_L = 0.6
TITLE_TOP = 0.35
BODY_TOP = 1.6

MAX_TABLE_ROWS = 14  # header + data rows cap per slide (keeps fonts legible)


class PowerPointReporter:
    """Generates a structural-analysis presentation deck (.pptx)."""

    def __init__(self):
        self.presentation: Optional[Presentation] = None

    # ------------------------------------------------------------------ #
    # Public API
    # ------------------------------------------------------------------ #

    def generate_presentation(
        self,
        file_path: str,
        project_title: str,
        engineer_name: str,
        summary_text: str,
        member_df: pd.DataFrame,
        reactions_df: pd.DataFrame,
        diagram_paths: Optional[List[str]] = None,
        design_standards: Optional[List[str]] = None,
        assumptions: Optional[List[str]] = None,
    ) -> str:
        """Builds and saves the full presentation deck; returns the file path."""
        self.presentation = Presentation()
        self.presentation.slide_width = Inches(SLIDE_W)
        self.presentation.slide_height = Inches(SLIDE_H)

        self._add_title_slide(project_title, engineer_name)
        self._add_bullets_slide("Project Assumptions", assumptions or self._default_assumptions())
        self._add_bullets_slide("Design Standards & Codes", design_standards or self._default_standards())
        self._add_summary_slide(summary_text)
        self._add_table_slide("Governing Member Forces", self._extract_governing_forces(member_df))
        self._add_table_slide("Support Reactions", reactions_df)

        for path in (diagram_paths or []):
            if path and os.path.isfile(path):
                self._add_image_slide(path)

        self._add_closing_slide()

        os.makedirs(os.path.dirname(os.path.abspath(file_path)) or ".", exist_ok=True)
        self.presentation.save(file_path)
        logger.info("Presentation saved to %s", file_path)
        return file_path

    # ------------------------------------------------------------------ #
    # Slide scaffolding
    # ------------------------------------------------------------------ #

    def _blank_slide(self):
        """Returns a new slide using the blank layout."""
        return self.presentation.slides.add_slide(self.presentation.slide_layouts[6])

    def _content_slide(self, title: str):
        """Returns a blank slide with a styled title + steel-blue accent line."""
        slide = self._blank_slide()

        box = slide.shapes.add_textbox(
            Inches(MARGIN_L), Inches(TITLE_TOP),
            Inches(SLIDE_W - 2 * MARGIN_L), Inches(0.85),
        )
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        run.font.size = Pt(28)
        run.font.bold = True
        run.font.name = "Calibri"
        run.font.color.rgb = NAVY

        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(MARGIN_L), Inches(1.32),
            Inches(SLIDE_W - 2 * MARGIN_L), Inches(0.045),
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = STEEL_BLUE
        accent.line.fill.background()
        return slide

    # ------------------------------------------------------------------ #
    # Individual slide builders
    # ------------------------------------------------------------------ #

    def _add_title_slide(self, project_title: str, engineer_name: str) -> None:
        slide = self._blank_slide()
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = NAVY

        box = slide.shapes.add_textbox(
            Inches(1.0), Inches(2.3), Inches(SLIDE_W - 2.0), Inches(1.6),
        )
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = project_title
        run.font.size = Pt(40)
        run.font.bold = True
        run.font.name = "Calibri"
        run.font.color.rgb = WHITE

        sub = slide.shapes.add_textbox(
            Inches(1.0), Inches(4.0), Inches(SLIDE_W - 2.0), Inches(1.2),
        )
        stf = sub.text_frame
        stf.word_wrap = True
        sp = stf.paragraphs[0]
        sp.alignment = PP_ALIGN.CENTER
        srun = sp.add_run()
        srun.text = "Structural Analysis Summary Presentation"
        srun.font.size = Pt(18)
        srun.font.name = "Calibri"
        srun.font.color.rgb = LIGHT_BLUE

        sp2 = stf.add_paragraph()
        sp2.alignment = PP_ALIGN.CENTER
        srun2 = sp2.add_run()
        srun2.text = (
            f"Prepared by: {engineer_name}  |  "
            f"{datetime.now().strftime('%Y-%m-%d %H:%M')}"
        )
        srun2.font.size = Pt(14)
        srun2.font.name = "Calibri"
        srun2.font.color.rgb = LIGHT_BLUE

        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(3.95), Inches(4.333), Inches(0.05),
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = STEEL_BLUE
        accent.line.fill.background()

    def _add_bullets_slide(self, title: str, bullets: List[str]) -> None:
        slide = self._content_slide(title)
        box = slide.shapes.add_textbox(
            Inches(MARGIN_L), Inches(BODY_TOP),
            Inches(SLIDE_W - 2 * MARGIN_L), Inches(SLIDE_H - BODY_TOP - 0.6),
        )
        tf = box.text_frame
        tf.word_wrap = True
        first = True
        for item in bullets:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            run = p.add_run()
            run.text = f"•  {item}"
            run.font.size = Pt(16)
            run.font.name = "Calibri"
            run.font.color.rgb = DARK_TEXT
            p.space_after = Pt(10)

    def _add_summary_slide(self, summary_text: str) -> None:
        slide = self._content_slide("Executive Summary")
        box = slide.shapes.add_textbox(
            Inches(MARGIN_L), Inches(BODY_TOP),
            Inches(SLIDE_W - 2 * MARGIN_L), Inches(SLIDE_H - BODY_TOP - 0.8),
        )
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = summary_text or "No summary provided."
        run.font.size = Pt(15)
        run.font.name = "Calibri"
        run.font.color.rgb = DARK_TEXT
        p.line_spacing = 1.25

    def _add_table_slide(self, title: str, df: pd.DataFrame) -> None:
        slide = self._content_slide(title)

        if df is None or df.empty:
            box = slide.shapes.add_textbox(
                Inches(MARGIN_L), Inches(BODY_TOP),
                Inches(SLIDE_W - 2 * MARGIN_L), Inches(1.0),
            )
            tf = box.text_frame
            run = tf.paragraphs[0].add_run()
            run.text = "No data available for this section."
            run.font.size = Pt(14)
            run.font.italic = True
            run.font.name = "Calibri"
            run.font.color.rgb = GREY
            return

        display = df.head(MAX_TABLE_ROWS - 1).reset_index(drop=True)
        n_rows = len(display) + 1  # + header row
        n_cols = len(display.columns)
        table_w = SLIDE_W - 2 * MARGIN_L
        row_h = 0.4 if n_rows <= 10 else 0.32

        shape = slide.shapes.add_table(
            n_rows, n_cols,
            Inches(MARGIN_L), Inches(BODY_TOP),
            Inches(table_w), Inches(row_h * n_rows),
        )
        table = shape.table
        table.first_row = True
        table.horz_banding = False

        col_w = Inches(table_w) // n_cols
        for c in range(n_cols):
            table.columns[c].width = col_w
        for r in range(n_rows):
            table.rows[r].height = Inches(row_h)

        # Header row
        for c, col_name in enumerate(display.columns):
            self._set_cell(
                table.cell(0, c), str(col_name).replace("_", " "),
                bold=True, color=WHITE, fill=STEEL_BLUE, size=12,
            )

        # Data rows
        for r, (_, row) in enumerate(display.iterrows(), start=1):
            fill = LIGHT_BLUE if r % 2 == 0 else WHITE
            for c, col_name in enumerate(display.columns):
                self._set_cell(
                    table.cell(r, c), self._fmt(row[col_name]),
                    bold=False, color=DARK_TEXT, fill=fill,
                    size=11 if n_rows <= 10 else 10,
                )

        if len(df) > len(display):
            note = slide.shapes.add_textbox(
                Inches(MARGIN_L), Inches(BODY_TOP + row_h * n_rows + 0.15),
                Inches(table_w), Inches(0.4),
            )
            nrun = note.text_frame.paragraphs[0].add_run()
            nrun.text = f"(showing first {len(display)} of {len(df)} rows)"
            nrun.font.size = Pt(11)
            nrun.font.italic = True
            nrun.font.name = "Calibri"
            nrun.font.color.rgb = GREY

    def _add_image_slide(self, image_path: str) -> None:
        base = os.path.basename(image_path).upper()
        if "SFD" in base:
            title = "Shear Force Diagram (SFD)"
        elif "BMD" in base:
            title = "Bending Moment Diagram (BMD)"
        else:
            title = os.path.splitext(os.path.basename(image_path))[0].replace("_", " ")

        slide = self._content_slide(title)
        max_w = SLIDE_W - 2 * MARGIN_L - 1.0
        max_h = SLIDE_H - BODY_TOP - 0.7

        pic = slide.shapes.add_picture(
            image_path, Inches(MARGIN_L), Inches(BODY_TOP), width=Inches(max_w),
        )
        if pic.height > Inches(max_h):
            ratio = Inches(max_h) / pic.height
            pic.width = int(pic.width * ratio)
            pic.height = int(pic.height * ratio)
        pic.left = int((Inches(SLIDE_W) - pic.width) / 2)
        pic.top = int(BODY_TOP + (Inches(max_h) - pic.height) / 2)

    def _add_closing_slide(self) -> None:
        slide = self._blank_slide()
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = NAVY

        box = slide.shapes.add_textbox(
            Inches(1.0), Inches(3.0), Inches(SLIDE_W - 2.0), Inches(1.5),
        )
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = "End of Presentation"
        run.font.size = Pt(30)
        run.font.bold = True
        run.font.name = "Calibri"
        run.font.color.rgb = WHITE

        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = (
            "Generated by the Structural Multi-App Agent — results derived "
            "from Autodesk Robot Structural Analysis Professional. "
            f"{datetime.now().strftime('%Y-%m-%d %H:%M')}"
        )
        run2.font.size = Pt(12)
        run2.font.name = "Calibri"
        run2.font.color.rgb = LIGHT_BLUE

    # ------------------------------------------------------------------ #
    # Helpers
    # ------------------------------------------------------------------ #

    @staticmethod
    def _set_cell(cell, text: str, bold: bool, color, fill, size: int) -> None:
        """Writes styled text into a table cell."""
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Inches(0.06)
        cell.margin_right = Inches(0.06)
        cell.margin_top = Inches(0.02)
        cell.margin_bottom = Inches(0.02)

        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.bold = bold
        run.font.size = Pt(size)
        run.font.name = "Calibri"
        run.font.color.rgb = color

    @staticmethod
    def _fmt(value) -> str:
        """Formats cell values: floats to 2 decimals, everything else str()."""
        if isinstance(value, float):
            return f"{value:.2f}"
        return str(value)

    # ------------------------------------------------------------------ #
    # Default content
    # ------------------------------------------------------------------ #

    @staticmethod
    def _default_assumptions() -> List[str]:
        return [
            "All members are modeled as linear-elastic prismatic beam elements.",
            "Self-weight of structural steel is included via material density unless stated otherwise.",
            "Supports are idealized as indicated in the model (fixed / pinned / roller) with no soil-structure interaction considered.",
            "Loads are applied as static, non-dynamic actions; seismic and wind dynamic amplification are outside the scope of this presentation unless separately noted.",
            "Second-order (P-Delta) effects are neglected unless explicitly enabled in the analysis case.",
        ]

    @staticmethod
    def _default_standards() -> List[str]:
        return [
            "Eurocode 3 (EN 1993-1-1) — Design of steel structures, General rules.",
            "Eurocode 1 (EN 1991-1-1) — Actions on structures, Densities, self-weight, imposed loads.",
            "Eurocode 0 (EN 1990) — Basis of structural design, load combination factors.",
            "Autodesk Robot Structural Analysis Professional — Finite Element solver, linear static analysis.",
        ]

    @staticmethod
    def _extract_governing_forces(member_df: pd.DataFrame) -> pd.DataFrame:
        """Reduces a full member-forces export to one governing (max |MY|) row per bar."""
        if member_df is None or member_df.empty:
            return pd.DataFrame(columns=["Bar_ID", "Position_m", "FX_kN", "FZ_kN", "MY_kNm"])

        if "MY_kNm" not in member_df.columns:
            logger.warning("member_df missing 'MY_kNm' column; returning empty governing forces.")
            return pd.DataFrame(columns=["Bar_ID", "Position_m", "FX_kN", "FZ_kN", "MY_kNm"])

        df = member_df.copy()
        if "Bar_ID" not in df.columns:
            return df

        df["_abs_my"] = df["MY_kNm"].abs()
        idx = df.groupby("Bar_ID")["_abs_my"].idxmax()
        governing = df.loc[idx].drop(columns="_abs_my").sort_values("Bar_ID").reset_index(drop=True)
        return governing



